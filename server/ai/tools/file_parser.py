"""
File parsers — extract text from PDF, PPTX, TXT uploads.

Each parser implements the same interface: bytes in → str out.
No vendor-specific dependencies leak into the rest of the system.
"""

from __future__ import annotations
import io
from typing import Protocol, runtime_checkable


@runtime_checkable
class FileParser(Protocol):
    """Any object with a `parse(raw_bytes, filename) -> str` method."""
    def parse(self, raw_bytes: bytes, filename: str) -> str: ...


class PDFParser:
    """Extract text from PDF files using pypdf."""

    def parse(self, raw_bytes: bytes, filename: str) -> str:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(raw_bytes))
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"--- Page {i + 1} ---\n{text.strip()}")
        if not pages:
            raise ValueError(f"Could not extract text from PDF '{filename}'. It may be image-based.")
        return "\n\n".join(pages)


class PPTXParser:
    """Extract text from PowerPoint files using python-pptx."""

    def parse(self, raw_bytes: bytes, filename: str) -> str:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw_bytes))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
        if not slides:
            raise ValueError(f"Could not extract text from PPTX '{filename}'.")
        return "\n\n".join(slides)


class TXTParser:
    """Parse plain text files."""

    def parse(self, raw_bytes: bytes, filename: str) -> str:
        # Try UTF-8 first, fall back to latin-1
        try:
            text = raw_bytes.decode("utf-8")
        except UnicodeDecodeError:
            text = raw_bytes.decode("latin-1")
        if not text.strip():
            raise ValueError(f"File '{filename}' is empty.")
        return text.strip()


# ──────────────────────────────────────────────────────────────
# Registry — maps file extensions to parser instances
# ──────────────────────────────────────────────────────────────

_PARSER_MAP: dict[str, FileParser] = {
    ".pdf": PDFParser(),
    ".pptx": PPTXParser(),
    ".ppt": PPTXParser(),
    ".txt": TXTParser(),
    ".md": TXTParser(),
    ".text": TXTParser(),
}

SUPPORTED_EXTENSIONS = set(_PARSER_MAP.keys())


def parse_file(raw_bytes: bytes, filename: str) -> str:
    """Parse a file based on its extension.

    Args:
        raw_bytes: Raw file content
        filename: Original filename (used to determine parser)

    Returns:
        Extracted text content

    Raises:
        ValueError: If file type is unsupported or content can't be extracted
    """
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    parser = _PARSER_MAP.get(ext)
    if not parser:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"Unsupported file type '{ext}'. Supported: {supported}")
    return parser.parse(raw_bytes, filename)
